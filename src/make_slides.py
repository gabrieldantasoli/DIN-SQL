"""
make_slides.py — Renderiza um deck (JSON) em .pptx tematico e ludico, com as
figuras embutidas e notas do apresentador.

JSON esperado:
{ "title": "...", "subtitle": "...", "author": "...",
  "slides": [ {"title","bullets":[...],"big":"...","image":"figN.png","notes":"..."} ] }

Uso: python src/make_slides.py deck.json APRESENTACAO.pptx
"""

import json
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

FIGDIR = "results/figures"
# paleta ludica/profissional
ACCENT = RGBColor(0x2D, 0x6C, 0xDF)   # azul
ACCENT2 = RGBColor(0x11, 0xA8, 0x8A)  # teal
DARK = RGBColor(0x1A, 0x22, 0x33)
GRAY = RGBColor(0x5B, 0x66, 0x77)
LIGHT = RGBColor(0xEF, 0xF4, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)


def _fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _bg(slide, color):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    _fill(r, color)
    slide.shapes._spTree.remove(r._element); slide.shapes._spTree.insert(2, r._element)
    return r


def _text(slide, left, top, width, height, text, size, color, bold=False,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
    return tb


def title_slide(prs, deck):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, DARK)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), EMU_W, Inches(0.12))
    _fill(band, ACCENT2)
    _text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.5),
          deck.get("title", ""), 40, WHITE, bold=True)
    if deck.get("subtitle"):
        _text(s, Inches(0.9), Inches(2.75), Inches(11.5), Inches(1.2),
              deck["subtitle"], 22, RGBColor(0xCF, 0xDB, 0xEC))
    if deck.get("author"):
        _text(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(1.0),
              deck["author"], 16, RGBColor(0x9F, 0xB0, 0xC6))


def content_slide(prs, sl, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    # barra de titulo
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, Inches(1.1))
    _fill(bar, ACCENT)
    _text(s, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.9),
          sl.get("title", ""), 28, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)

    img = sl.get("image", "") or ""
    img_path = os.path.join(FIGDIR, img) if img else ""
    has_img = bool(img) and os.path.exists(img_path)

    body_w = Inches(6.4) if has_img else Inches(12.3)
    top = Inches(1.45)

    # destaque "big"
    if sl.get("big"):
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, body_w, Inches(1.0))
        _fill(box, LIGHT); box.line.color.rgb = ACCENT2; box.line.width = Pt(1.5)
        tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = sl["big"]
        r.font.size = Pt(19); r.font.bold = True; r.font.color.rgb = ACCENT
        top = Inches(2.65)

    # bullets
    bullets = sl.get("bullets", []) or []
    if bullets:
        tb = s.shapes.add_textbox(Inches(0.5), top, body_w, Inches(4.6))
        tf = tb.text_frame; tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(10)
            lead = "" if (b[:1] and not b[:1].isalnum()) else "•  "
            r = p.add_run(); r.text = lead + b
            r.font.size = Pt(18); r.font.color.rgb = DARK; r.font.name = "Calibri"

    # imagem (direita), preservando proporcao
    if has_img:
        try:
            from PIL import Image
            iw, ih = Image.open(img_path).size
            ar = iw / ih
        except Exception:
            ar = 1.6
        maxw, maxh = Inches(6.0), Inches(5.0)
        w = maxw; h = int(w / ar)
        if h > maxh:
            h = maxh; w = int(h * ar)
        left = Inches(13.333) - Inches(0.45) - w
        s.shapes.add_picture(img_path, left, Inches(1.7), width=w, height=h)

    # rodape
    _text(s, Inches(0.5), Inches(6.95), Inches(9), Inches(0.4),
          "DIN-SQL · reprodução de baixo custo · COPIN/UFCG", 10, GRAY)
    _text(s, Inches(12.2), Inches(6.95), Inches(0.9), Inches(0.4),
          f"{idx}/{total}", 10, GRAY, align=PP_ALIGN.RIGHT)

    # notas do apresentador
    if sl.get("notes"):
        s.notes_slide.notes_text_frame.text = sl["notes"]


def main():
    deck = json.load(open(sys.argv[1]))
    out = sys.argv[2]
    prs = Presentation()
    prs.slide_width = EMU_W; prs.slide_height = EMU_H
    title_slide(prs, deck)
    slides = deck.get("slides", [])
    for i, sl in enumerate(slides, 1):
        content_slide(prs, sl, i, len(slides))
    prs.save(out)
    n_img = sum(1 for sl in slides if (sl.get("image") and
               os.path.exists(os.path.join(FIGDIR, sl["image"]))))
    print(f"{out} gerado — {len(slides)+1} slides ({n_img} com figura)")


if __name__ == "__main__":
    main()
